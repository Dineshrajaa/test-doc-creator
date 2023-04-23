import os
import argparse
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx2pdf import convert
import subprocess
import json
from pkg_resources import parse_version as LooseVersion
from PIL import Image, ImageDraw, ImageFont
import xml.etree.ElementTree as ET

parser = argparse.ArgumentParser(description="Generate files with embedded text.")
parser.add_argument("contractNumber", type=str, help="the contract or claim number")
parser.add_argument(
    "type", type=str, choices=["Contract", "Claim"], help="the type of the document"
)
args = parser.parse_args()

contractNumber = args.contractNumber
docType = args.type

lorem_ipsum = (
    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Donec at nulla elit. Nulla pharetra purus turpis, in volutpat justo aliquet eget. Suspendisse malesuada nulla eget tellus hendrerit fringilla. Duis tincidunt elit eget elit molestie viverra. Quisque lobortis consequat lorem, sit amet dictum justo malesuada id. Ut fermentum bibendum massa, non finibus ex feugiat id. Praesent euismod orci vitae dictum tristique. Aenean euismod euismod massa, in blandit augue vehicula sit amet. Pellentesque habitant morbi tristique senectus et netus et malesuada fames ac turpis egestas. Nam pulvinar, tortor sed porttitor placerat, sem mauris fringilla nibh, a varius augue massa vel sapien. Fusce condimentum ornare tortor, eget gravida velit. Quisque venenatis massa eu enim tincidunt vestibulum. Donec id tincidunt orci. Curabitur egestas nulla vel ultrices luctus. Sed quis tortor non metus eleifend accumsan. Nulla facilisi. \n"
    * 50
)

# Embed the contract/claim number in the middle of the Lorem Ipsum text
with_number = f"{lorem_ipsum}test data is here {contractNumber}"
# Create a directory with the contract/claim number as its name
directory = os.path.join(os.getcwd(), f"generated/{contractNumber}")
if not os.path.exists(directory):
    os.makedirs(directory)

# Create the specified files in the directory
filetypes = [
    "pdf",
    "txt",
    "xls",
    "xlsx",
    "xlsb",
    "xlsm",
    "html",
    "htm",
    "csv",
    "rtf",
    "rar",
    "ppt",
    "pptx",
    "doc",
    "xml",
    "json",
    "png",
]
# JSON and XML contents
jsonXMLData = {
    'name': 'John Doe',
    'age': 30,
    'city': 'New York',
    'address': {
        'street': '123 Main St',
        'city': 'New York',
        'state': 'NY',
        'zip': '10001',
        'country': 'USA'
    },
    'phone': [
        {'type': 'home', 'number': '555-1234'},
        {'type': 'work', 'number': '555-5678'}
    ],
    'email': ['john.doe@example.com', 'johndoe@gmail.com'],
    'preferences': {
        'color': 'blue',
        'food': ['pizza', 'sushi', 'tacos'],
        'hobbies': ['reading', 'hiking', 'playing guitar']
    },
    'order': {
        'type': f"{docType}",
        'id': f"{contractNumber}"
    }
}

for filetype in filetypes:
    filename = os.path.join(directory, f"{docType}.{filetype}")
    with open(filename, "w") as file:
        file.write(with_number)
    if filetype.startswith("xls"):
        try:
            sheetname = f"{docType}"
            df = pd.DataFrame({"text": [with_number]})
            excel_filename = os.path.join(directory, f"{docType}.xlsx")
            df.to_excel(excel_filename, sheet_name=sheetname, index=False)
        except ImportError:
            print(f"Warning: Pandas not found, cannot generate {filetype} file")
    elif filetype == "pdf":
        document = Document()
        documentFileName = os.path.join(directory, f"{docType}.docx")
        for _ in range(5):
            document.add_paragraph(with_number, style="Normal")
        document.save(documentFileName)
        if os.name == "posix":
            subprocess.run(["unoconv", "-f", "pdf", documentFileName])
        else:
            convert(documentFileName, os.path.join(directory, f"{docType}.pdf"))
    elif filetype.startswith("ppt"):
        try:
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = f"{docType} Presentation"
            minimumppttext = (
                "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Donec at nulla elit.\n"
                * 2
            )
            subtitle.text = f"{minimumppttext}test data{contractNumber}"
            ppt_filename = os.path.join(
                directory, f"{docType}.pptx" if filetype == "pptx" else f"{docType}.ppt"
            )
            prs.save(ppt_filename)
        except ImportError:
            print(f"Warning: python-pptx not found, cannot generate {filetype} file")
    elif filetype == "xml":
        root = ET.Element('jsonXMLData')
        for key, value in jsonXMLData.items():
            param = ET.SubElement(root, 'param')
            param.set('name', key)
            param.text = str(value)
        tree = ET.ElementTree(root)
        xml_path = os.path.join(directory, f"{docType}.xml")
        tree.write(xml_path)
    elif filetype == "json":
        json_path = os.path.join(directory, f"{docType}.json")
        with open(json_path, 'w') as json_file:
            json.dump(jsonXMLData, json_file)
    elif filetype == "png":
        background_color = (255, 255, 255)
        font_size = 60
        font_type = "./roboto.ttf"
        font = ImageFont.truetype(font_type, font_size)
        text_color = (255, 255, 255, 255)
        image = Image.open("bgimage.jpg")
        draw = ImageDraw.Draw(image)
        width, height = image.size
        text = f"{docType}-{contractNumber}"
        text_width, text_height = draw.textsize(text, font)
        # Calculate the position to draw the text (centered on the image)
        text_x = (width - text_width) / 2
        text_y = (height - text_height) / 2
        draw.text((text_x, text_y), text, fill=text_color, font=font)
        png_path = os.path.join(directory, f"{docType}.png")
        image.save(png_path, format="PNG")
        print(f"Use this file path with image generator script {png_path}")

