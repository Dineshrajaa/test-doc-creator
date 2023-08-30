import os
import pandas as pd
import xlwt  # For xls format
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx2pdf import convert
import subprocess
import json
from PIL import Image, ImageDraw, ImageFont
import xml.etree.ElementTree as ET
from email.message import EmailMessage
# from msglib import Message

# Read the configuration excel file
config_df = pd.read_excel("config.xlsx", sheet_name="Doc Types")
config_df.dropna(subset=["File Types"], inplace=True)

# Convert float values in "Is Required" column to boolean
config_df["Is Required"] = config_df["Is Required"].astype(int).astype(bool)

# Get the required file types based on the "Is Required" column
required_filetypes = config_df[config_df["Is Required"]]["File Types"].tolist()

# Read IDs and ID Types from Sheet 2
id_data = pd.read_excel("config.xlsx", sheet_name="Test IDs")

summary_data = []


def generate_excel_files(doc_type, with_number, directory):
    try:
        sheetname = f"{doc_type}"
        df = pd.DataFrame({"text": [with_number]})

        # Generate xlsx file
        xlsx_filename = os.path.join(directory, f"{doc_type}.xlsx")
        df.to_excel(xlsx_filename, sheet_name=sheetname, index=False)

        # Generate xls file (requires xlwt)
        xls_filename = os.path.join(directory, f"{doc_type}.xls")
        workbook = xlwt.Workbook()
        sheet = workbook.add_sheet(sheetname)
        sheet.write(0, 0, "text")
        sheet.write(1, 0, with_number)
        workbook.save(xls_filename)

    except ImportError:
        print("Warning: Some required libraries not found, cannot generate Excel files")


def generate_image_base(doc_type, doc_id):
    font_size = 60
    font_type = "./roboto.ttf"
    font = ImageFont.truetype(font_type, font_size)
    text_color = (255, 255, 255, 255)
    image = Image.open("bgimage.jpg")
    draw = ImageDraw.Draw(image)
    width, height = image.size
    text = f"{doc_type}-{doc_id}"
    text_width, text_height = draw.textsize(text, font)
    # Calculate the position to draw the text (centered on the image)
    text_x = (width - text_width) / 2
    text_y = (height - text_height) / 2
    draw.text((text_x, text_y), text, fill=text_color, font=font)
    return image


def generate_html_file(doc_type, doc_id):
    html_content = f"""<!DOCTYPE html>
                        <html>
                            <head>
                                <title>{doc_type}</title>
                            </head>
                            <body>
                                <h1>Hello, Tester!</h1>
                                <p>This is having the test data of {doc_type} {doc_id}</p>
                            </body>
                            </html>
"""
    html_types = ["htm", "html"]
    for html_type in html_types:
        filename = os.path.join(directory, f"{doc_type}.{html_type}")
        with open(filename, "w") as file:
            file.write(html_content)


def generate_eml_file(doc_type, doc_id):
    msg = EmailMessage()
    msg["From"] = "thangameena.nagarajan@allianz.com"
    msg["To"] = "thangameena.nagarajan@allianz.com"
    msg["Subject"] = f"Automatically created to test {doc_type} {doc_id}"
    msg.set_content(f"{lorem_ipsum} this is the test data of {doc_type} {doc_id}")
    filename = os.path.join(directory, f"{doc_type}.eml")
    with open(filename, "w") as file:
        file.write(msg.as_string())

# TBD: Enable msg (outlook) file generation with proper package
# def generate_msg_file(doc_type, doc_id):
#     msg = Message()
#     msg.sender = "thangameena.nagarajan@allianz.com"
#     msg.recipients = ["thangameena.nagarajan@allianz.com"]
#     msg.subject = f"Automatically created to test {doc_type} {doc_id}"
#     msg.body = f"{lorem_ipsum} this is the test data of {doc_type} {doc_id}"
#     filename = os.path.join(directory, f"{doc_type}.msg")
#     msg.save(filename)

# Iterate through the rows of the configuration file
for _, id_row in id_data.iterrows():
    doc_type = id_row["ID Type"]
    doc_id = id_row["ID"]

    # Create a directory with the doc_id as its name
    directory = os.path.join(os.getcwd(), f"generated/{doc_id}")
    if not os.path.exists(directory):
        os.makedirs(directory)

    lorem_ipsum = "Lorem ipsum dolor sit amet. \n" * 10
    # Embed the contract/claim number in the middle of the Lorem Ipsum text
    with_number = f"{lorem_ipsum}test data is here {doc_id}"
    # JSON and XML contents
    jsonXMLData = {
        "name": "John Doe",
        "age": 30,
        "city": "New York",
        "address": {
            "street": "123 Main St",
            "city": "New York",
            "state": "NY",
            "zip": "10001",
            "country": "USA",
        },
        "phone": [
            {"type": "home", "number": "555-1234"},
            {"type": "work", "number": "555-5678"},
        ],
        "email": ["john.doe@example.com", "johndoe@gmail.com"],
        "preferences": {
            "color": "blue",
            "food": ["pizza", "sushi", "tacos"],
            "hobbies": ["reading", "hiking", "playing guitar"],
        },
        "order": {"type": f"{doc_type}", "id": f"{doc_id}"},
    }
    # Generate files only for required file types
    for filetype in required_filetypes:
        if filetype.startswith("xls"):
            generate_excel_files(doc_type, with_number, directory)
        elif filetype == "pdf":
            document = Document()
            documentFileName = os.path.join(directory, f"{doc_type}.docx")
            for _ in range(5):
                document.add_paragraph(lorem_ipsum, style="Normal")
            document.add_paragraph(with_number, style="Normal")
            document.save(documentFileName)
            if os.name == "posix":
                subprocess.run(["unoconv", "-f", "pdf", documentFileName])
            else:
                convert(documentFileName, os.path.join(directory, f"{doc_type}.pdf"))
        elif filetype.startswith("ppt"):
            try:
                prs = Presentation()
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = f"{doc_type} Presentation"
                minimumppttext = (
                    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Donec at nulla elit.\n"
                    * 2
                )
                subtitle.text = f"{minimumppttext}test data{doc_id}"
                ppt_filename = os.path.join(
                    directory,
                    f"{doc_type}.pptx" if filetype == "pptx" else f"{doc_type}.ppt",
                )
                prs.save(ppt_filename)
            except ImportError:
                print(
                    f"Warning: python-pptx not found, cannot generate {filetype} file"
                )
        elif filetype == "xml":
            root = ET.Element("jsonXMLData")
            for key, value in jsonXMLData.items():
                param = ET.SubElement(root, "param")
                param.set("name", key)
                param.text = str(value)
            tree = ET.ElementTree(root)
            xml_path = os.path.join(directory, f"{doc_type}.xml")
            tree.write(xml_path)
        elif filetype == "json":
            json_path = os.path.join(directory, f"{doc_type}.json")
            with open(json_path, "w") as json_file:
                json.dump(jsonXMLData, json_file)
        elif filetype == "png":
            image = generate_image_base(doc_type, doc_id)
            png_path = os.path.join(directory, f"{doc_type}.png")
            image.save(png_path, format="PNG")
        elif filetype == "bmp":
            image = generate_image_base(doc_type, doc_id)
            png_path = os.path.join(directory, f"{doc_type}.bmp")
            image.save(png_path, format="BMP")
        elif filetype == "jpeg":
            image = generate_image_base(doc_type, doc_id)
            png_path = os.path.join(directory, f"{doc_type}.jpeg")
            image.save(png_path, format="JPEG")
        elif filetype == "tiff":
            image = generate_image_base(doc_type, doc_id)
            png_path = os.path.join(directory, f"{doc_type}.tiff")
            image.save(png_path, format="TIFF")
        elif filetype == "gif":
            image = generate_image_base(doc_type, doc_id)
            png_path = os.path.join(directory, f"{doc_type}.gif")
            image.save(png_path, format="GIF")
        elif filetype == "htm":
            generate_html_file(doc_type, doc_id)
        elif filetype == "eml":
            generate_eml_file(doc_type, doc_id)
        # elif filetype == "msg":
        #     generate_msg_file(doc_type, doc_id)
        else:
            filename = os.path.join(directory, f"{doc_type}.{filetype}")
            with open(filename, "w") as file:
                file.write(with_number)  # You need to define "with_number"

    # Create summary data
    for filetype in required_filetypes:
        summary_data.append(
            {
                "File type": filetype,
                "ID": doc_id,
                "ID Type": doc_type,
                "File Name": f"{doc_type}.{filetype}",
                "Saved path": os.path.join(directory, f"{doc_type}.{filetype}"),
            }
        )

    # Create a summary Excel file
    summary_df = pd.DataFrame(summary_data)
    summary_excel_path = os.path.join(os.getcwd(), f"summary.xlsx")
    summary_df.to_excel(summary_excel_path, index=False)
